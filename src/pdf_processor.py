import os
import re
import cv2
import numpy as np
from PIL import Image
from loguru import logger

try:
    import pytesseract
except ImportError:
    pytesseract = None
try:
    import easyocr
except ImportError:
    easyocr = None

# ---- 파일 포맷별 프로세서 ----
class BaseFileProcessor:
    def extract_text(self, file_path: str) -> str:
        raise NotImplementedError

    def extract_chunks(self, file_path: str, department: str = None) -> list:
        raise NotImplementedError

class PDFProcessor(BaseFileProcessor):
    def extract_text(self, file_path: str) -> str:
        import pdfplumber
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text.append(page_text)
        return "\n".join(text)

    def extract_chunks(self, file_path: str, department: str = None) -> list:
        import pdfplumber
        chunks = []
        filename = os.path.basename(file_path)
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                # 헤더/푸터/목차 등 제거 (간단 예시)
                lines = [l for l in page_text.split('\n') if not re.match(r'^(목차|Page|페이지|Copyright|All rights reserved)', l.strip())]
                page_text_clean = '\n'.join(lines)
                # 문단 단위로 쪼개기 (빈 줄 기준)
                paragraphs = [p.strip() for p in re.split(r'\n\s*\n', page_text_clean) if p.strip()]
                for para in paragraphs:
                    for i in range(0, len(para), 1000):
                        chunk = para[i:i+1000]
                        meta = {
                            "type": "pdf",
                            "filename": filename,
                            "page": page_num,
                            "department": department
                        }
                        chunks.append({"text": chunk, "metadata": meta})
        return chunks

class ExcelProcessor(BaseFileProcessor):
    def extract_text(self, file_path: str) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_blocks = []
        for sheet in wb.worksheets:
            try:
                rows = list(sheet.iter_rows(values_only=True))
            except Exception as e:
                logger.error(f"[ExcelProcessor] 시트 '{sheet.title}'에서 행을 읽는 중 오류: {e}")
                continue
            if not rows or all(row is None or all(cell is None for cell in (row or [])) for row in rows):
                continue
            header_row = rows[0] if len(rows) > 0 else None
            if header_row is None or all(cell is None or str(cell).strip() == "" for cell in (header_row or [])):
                continue
            header = [str(cell) if cell is not None else "" for cell in (header_row or [])]
            for row_idx, row in enumerate(rows[1:], start=2):
                if row is None or not isinstance(row, (list, tuple)) or all(cell is None for cell in (row or [])):
                    continue
                # row 길이가 header보다 짧으면 부족한 부분은 None으로 채움
                row_filled = list(row) + [None] * (len(header) - len(row))
                row_dict = {k: v for k, v in zip(header, row_filled)}
                if not any(v is not None and str(v).strip() != "" for v in row_dict.values()):
                    continue
                text = f"[{sheet.title}] " + ", ".join(f"{k}: {v}" for k, v in row_dict.items() if v is not None and str(v).strip() != "")
                text_blocks.append(text)
        return "\n".join(text_blocks)

    def extract_chunks(self, file_path: str, department: str = None) -> list:
        import openpyxl
        filename = os.path.basename(file_path)
        chunks = []
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            try:
                rows = list(sheet.iter_rows(values_only=True))
            except Exception as e:
                logger.error(f"[ExcelProcessor] 시트 '{sheet.title}'에서 행을 읽는 중 오류: {e}")
                continue
            if not rows or all(row is None or all(cell is None for cell in (row or [])) for row in rows):
                continue
            header_row = rows[0] if len(rows) > 0 else None
            if header_row is None or all(cell is None or str(cell).strip() == "" for cell in (header_row or [])):
                continue
            header = [str(cell) if cell is not None else "" for cell in (header_row or [])]
            for row_idx, row in enumerate(rows[1:], start=2):
                if row is None or not isinstance(row, (list, tuple)) or all(cell is None for cell in (row or [])):
                    continue
                row_filled = list(row) + [None] * (len(header) - len(row))
                row_dict = {k: v for k, v in zip(header, row_filled)}
                if not any(v is not None and str(v).strip() != "" for v in row_dict.values()):
                    continue
                # 항상 모든 칼럼을 text로 사용
                text = ", ".join(f"{k}: {v}" for k, v in row_dict.items() if v is not None and str(v).strip() != "")
                meta = {
                    "type": "excel",
                    "sheet": sheet.title,
                    "row": row_idx,
                    "filename": filename,
                    "department": department,
                }
                # 모든 칼럼을 메타데이터에 추가
                for k, v in row_dict.items():
                    if v is not None and str(v).strip() != "":
                        meta[k] = v
                if text.strip():
                    chunks.append({"text": text, "metadata": meta})
        return chunks

class PowerPointProcessor(BaseFileProcessor):
    def extract_text(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_chunks(self, file_path: str, department: str = None) -> list:
        from pptx import Presentation
        filename = os.path.basename(file_path)
        prs = Presentation(file_path)
        chunks = []
        for idx, slide in enumerate(prs.slides, 1):
            title = ""
            body = []
            notes = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if hasattr(shape, 'shape_type') and shape.shape_type == 1:
                        title = shape.text.strip()
                    else:
                        body.append(shape.text.strip())
            if hasattr(slide, "notes_slide") and hasattr(slide.notes_slide, "notes_text_frame"):
                notes = slide.notes_slide.notes_text_frame.text.strip()
            text = f"[슬라이드 {idx}] {title}\n" + " ".join(body)
            if notes:
                text += f"\n[노트] {notes}"
            meta = {
                "type": "pptx",
                "slide": idx,
                "title": title,
                "filename": filename,
                "department": department
            }
            chunks.append({"text": text, "metadata": meta})
        return chunks

class ImageProcessor(BaseFileProcessor):
    def extract_text(self, file_path: str) -> str:
        img = cv2.imread(file_path)
        img = preprocess_image_for_ocr(img)
        text = ""
        if easyocr:
            try:
                reader = easyocr.Reader(['ko', 'en'])
                text = "\n".join([r[1] for r in reader.readtext(img)])
            except Exception as e:
                print(f"[easyocr 오류] {e}")
        if (not text or not text.strip()) and pytesseract:
            try:
                from PIL import Image
                pil_img = Image.fromarray(img)
                text = pytesseract.image_to_string(pil_img, lang='kor+eng')
            except Exception as e:
                print(f"[pytesseract 오류] {e}")
        return text

    def extract_chunks(self, file_path: str, department: str = None) -> list:
        img = cv2.imread(file_path)
        img = preprocess_image_for_ocr(img)
        text = ""
        if easyocr:
            try:
                reader = easyocr.Reader(['ko', 'en'])
                text = "\n".join([r[1] for r in reader.readtext(img)])
            except Exception as e:
                print(f"[easyocr 오류] {e}")
        if (not text or not text.strip()) and pytesseract:
            try:
                from PIL import Image
                pil_img = Image.fromarray(img)
                text = pytesseract.image_to_string(pil_img, lang='kor+eng')
            except Exception as e:
                print(f"[pytesseract 오류] {e}")
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        filename = os.path.basename(file_path)
        chunks = []
        for para in paragraphs:
            meta = {
                "type": "image",
                "filename": filename,
                "department": department
            }
            chunks.append({"text": para, "metadata": meta})
        return chunks

def preprocess_image_for_ocr(img):
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    kernel = np.array([[0, -1, 0], [-1, 5,-1], [0, -1, 0]])
    sharp = cv2.filter2D(gray, -1, kernel)
    sharp = cv2.convertScaleAbs(sharp, alpha=1.5, beta=0)
    proc = cv2.adaptiveThreshold(sharp, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 10)
    proc = cv2.fastNlMeansDenoising(proc, None, 30, 7, 21)
    proc = cv2.resize(proc, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
    return proc

def get_processor(file_path: str) -> BaseFileProcessor:
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        return PDFProcessor()
    elif ext == ".docx":
        return WordProcessor()
    elif ext == ".xlsx":
        return ExcelProcessor()
    elif ext == ".pptx":
        return PowerPointProcessor()
    elif ext in [".jpg", ".jpeg", ".png"]:
        return ImageProcessor()
    else:
        raise ValueError("지원하지 않는 파일 형식입니다.")
